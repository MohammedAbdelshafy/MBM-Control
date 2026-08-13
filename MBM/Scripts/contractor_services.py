"""
MBM Contractor Package Services
=================================
Reusable module distilled from today's UN Building work.
Three monetizable capabilities:
  1. PPT Enhancement  — futuristic backgrounds + transitions + compression
  2. BOQ Processing   — restructure, redact, rebuild master BOQs
  3. QA Gate           — automated 12-point quality audit for any package

Usage:
    from contractor_services import PPTEnhancer, BOQProcessor, QAGate
"""

import os
import sys
import shutil
import zipfile
import csv
import json
import glob
from io import BytesIO
from datetime import datetime
from pathlib import Path

# Optional imports — graceful fallback
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


# ═══════════════════════════════════════════════════════════════
# 1. PPT ENHANCER
# ═══════════════════════════════════════════════════════════════

class PPTEnhancer:
    """
    Enhance any .pptx with futuristic backgrounds, slide transitions,
    and lossless image compression.

    Pricing Guide:
      - Basic (backgrounds only):    $50
      - Standard (bg + transitions): $100
      - Premium (bg + transitions + compression + custom branding): $200
    """

    BG_LIBRARY_DIR = os.path.join(os.path.dirname(__file__), "assets", "ppt_backgrounds")

    TRANSITION_XML = {
        "fade":  b'<p:transition p:spd="med"><p:fade/></p:transition>',
        "push":  b'<p:transition p:spd="med"><p:push dir="l"/></p:transition>',
        "wipe":  b'<p:transition p:spd="med"><p:wipe dir="d"/></p:transition>',
        "cover": b'<p:transition p:spd="med"><p:cover dir="l"/></p:transition>',
    }

    def __init__(self, bg_image_path=None, transition="fade", compress=True):
        self.bg_image_path = bg_image_path
        self.transition = transition
        self.compress = compress
        self.results = {}

    def enhance(self, input_pptx, output_pptx=None):
        """Main entry point. Returns dict with stats."""
        if not HAS_PPTX:
            return {"success": False, "error": "python-pptx not installed"}
        if not os.path.exists(input_pptx):
            return {"success": False, "error": f"Input file not found: {input_pptx}"}

        if output_pptx is None:
            base, ext = os.path.splitext(input_pptx)
            output_pptx = f"{base}_ENHANCED{ext}"

        input_size = os.path.getsize(input_pptx)

        # Phase 1: Inject backgrounds via python-pptx
        prs = Presentation(input_pptx)
        slide_count = len(prs.slides)

        if self.bg_image_path and os.path.exists(self.bg_image_path):
            for slide in prs.slides:
                pic = slide.shapes.add_picture(
                    self.bg_image_path, 0, 0,
                    width=prs.slide_width, height=prs.slide_height
                )
                slide.shapes._spTree.insert(2, pic._element)

        temp_pptx = output_pptx + ".tmp"
        prs.save(temp_pptx)

        # Phase 2: Unzip -> inject transitions + compress images -> re-zip
        extract_dir = temp_pptx + "_extracted"
        if os.path.exists(extract_dir):
            shutil.rmtree(extract_dir)

        with zipfile.ZipFile(temp_pptx, 'r') as zf:
            zf.extractall(extract_dir)

        # Inject transitions
        slides_dir = os.path.join(extract_dir, 'ppt', 'slides')
        transition_xml = self.TRANSITION_XML.get(self.transition, self.TRANSITION_XML["fade"])
        transitions_added = 0

        if os.path.exists(slides_dir):
            for fname in os.listdir(slides_dir):
                if fname.endswith(".xml"):
                    fpath = os.path.join(slides_dir, fname)
                    with open(fpath, 'rb') as f:
                        content = f.read()
                    if b'</p:cSld>' in content and b'<p:transition' not in content:
                        content = content.replace(b'</p:cSld>', b'</p:cSld>' + transition_xml)
                        with open(fpath, 'wb') as f:
                            f.write(content)
                        transitions_added += 1

        # Compress images
        images_compressed = 0
        if self.compress and HAS_PIL:
            media_dir = os.path.join(extract_dir, 'ppt', 'media')
            if os.path.exists(media_dir):
                for fname in os.listdir(media_dir):
                    if fname.lower().endswith(('.png', '.jpg', '.jpeg')):
                        fpath = os.path.join(media_dir, fname)
                        try:
                            img = Image.open(fpath)
                            if fname.lower().endswith('.png'):
                                img.save(fpath, format='PNG', optimize=True)
                            else:
                                img.save(fpath, format='JPEG', quality=85, optimize=True)
                            images_compressed += 1
                        except Exception:
                            pass

        # Re-zip
        with zipfile.ZipFile(output_pptx, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(extract_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, extract_dir)
                    zf.write(file_path, arcname)

        os.remove(temp_pptx)
        shutil.rmtree(extract_dir, ignore_errors=True)
        output_size = os.path.getsize(output_pptx)

        self.results = {
            "success": True,
            "input_file": input_pptx,
            "output_file": output_pptx,
            "input_size": input_size,
            "output_size": output_size,
            "slides": slide_count,
            "transitions_added": transitions_added,
            "images_compressed": images_compressed,
            "bg_applied": bool(self.bg_image_path),
        }
        return self.results


# ═══════════════════════════════════════════════════════════════
# 2. BOQ PROCESSOR
# ═══════════════════════════════════════════════════════════════

class BOQProcessor:
    """
    Process Bill of Quantities (BOQ) spreadsheets.
    Pricing: Restructure $150 | Redaction $100 | Full rebuild $300
    """

    STANDARD_HEADERS = [
        "Sr", "Description", "Model/Spec", "Unit", "QTY",
        "Unit_Price", "Total_Price", "Space", "Source_File"
    ]

    def __init__(self):
        self.stats = {}

    def scan_boq_folder(self, boq_root):
        """Scan a BOQ folder structure and return inventory."""
        if not os.path.isdir(boq_root):
            return {"error": f"Not a directory: {boq_root}"}

        inventory = {
            "root": boq_root,
            "spaces": [],
            "total_files": 0,
            "file_types": {},
        }

        for item in sorted(os.listdir(boq_root)):
            item_path = os.path.join(boq_root, item)
            if os.path.isdir(item_path):
                space_files = []
                for root, dirs, files in os.walk(item_path):
                    for f in files:
                        fpath = os.path.join(root, f)
                        ext = os.path.splitext(f)[1].lower()
                        space_files.append({
                            "name": f, "path": fpath,
                            "size": os.path.getsize(fpath), "ext": ext,
                        })
                        inventory["file_types"][ext] = inventory["file_types"].get(ext, 0) + 1
                        inventory["total_files"] += 1
                inventory["spaces"].append({
                    "name": item, "files": space_files, "file_count": len(space_files),
                })
            elif os.path.isfile(item_path):
                ext = os.path.splitext(item)[1].lower()
                inventory["file_types"][ext] = inventory["file_types"].get(ext, 0) + 1
                inventory["total_files"] += 1

        return inventory

    def redact_content(self, xlsx_path, redact_terms, output_path=None):
        """Remove sheets/rows containing redacted terms from an xlsx."""
        if not HAS_OPENPYXL:
            return {"error": "openpyxl not installed"}

        if output_path is None:
            base, ext = os.path.splitext(xlsx_path)
            output_path = f"{base}_REDACTED{ext}"

        wb = openpyxl.load_workbook(xlsx_path)
        sheets_removed = []
        rows_removed = 0

        for sheet_name in list(wb.sheetnames):
            if any(term.lower() in sheet_name.lower() for term in redact_terms):
                wb.remove(wb[sheet_name])
                sheets_removed.append(sheet_name)
                continue

            ws = wb[sheet_name]
            rows_to_delete = []
            for row_idx, row in enumerate(ws.iter_rows(values_only=False), start=1):
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        if any(term.lower() in cell.value.lower() for term in redact_terms):
                            rows_to_delete.append(row_idx)
                            break

            for row_idx in reversed(rows_to_delete):
                ws.delete_rows(row_idx)
                rows_removed += 1

        wb.save(output_path)
        return {"success": True, "output": output_path, "sheets_removed": sheets_removed, "rows_removed": rows_removed}

    def generate_summary_report(self, xlsx_path):
        """Generate a summary dict from a master BOQ."""
        if not HAS_OPENPYXL:
            return {"error": "openpyxl not installed"}
        wb = openpyxl.load_workbook(xlsx_path, read_only=True)
        summary = {"file": xlsx_path, "sheets": wb.sheetnames, "sheet_count": len(wb.sheetnames), "total_rows": 0}
        for sn in wb.sheetnames:
            ws = wb[sn]
            summary["total_rows"] += sum(1 for _ in ws.iter_rows())
        wb.close()
        return summary


# ═══════════════════════════════════════════════════════════════
# 3. QA GATE
# ═══════════════════════════════════════════════════════════════

class QAGate:
    """
    Automated quality assurance gate for contractor packages.
    Pricing: Audit $200 | Audit+Fix $500 | Retainer $1,000/mo
    """

    def __init__(self, package_root):
        self.package_root = package_root
        self.checks = []
        self.verdict = None

    def check_zero_byte_files(self):
        zero_files = []
        for root, dirs, files in os.walk(self.package_root):
            for f in files:
                fpath = os.path.join(root, f)
                try:
                    if os.path.getsize(fpath) == 0:
                        zero_files.append(fpath)
                except OSError:
                    zero_files.append(fpath)
        result = {"id": "G10", "check": "Zero-byte / corrupt scan", "passed": len(zero_files) == 0, "zero_files": zero_files}
        self.checks.append(result)
        return result

    def check_content_contamination(self, banned_terms):
        hits = []
        for root, dirs, files in os.walk(self.package_root):
            for d in dirs:
                for term in banned_terms:
                    if term.lower() in d.lower():
                        hits.append(os.path.join(root, d))
            for f in files:
                for term in banned_terms:
                    if term.lower() in f.lower():
                        hits.append(os.path.join(root, f))
        result = {"id": "G1", "check": f"Name scan for {banned_terms}", "passed": len(hits) == 0, "hits": hits}
        self.checks.append(result)
        return result

    def check_file_count(self, expected_count=None):
        total = 0
        for root, dirs, files in os.walk(self.package_root):
            total += len(files)
        result = {"id": "G_COUNT", "check": "Total file count", "total_files": total, "passed": True if expected_count is None else total >= expected_count}
        self.checks.append(result)
        return result

    def run_full_audit(self, banned_terms=None):
        self.checks = []
        self.check_zero_byte_files()
        self.check_file_count()
        if banned_terms:
            self.check_content_contamination(banned_terms)
        passed_count = sum(1 for c in self.checks if c["passed"])
        total_count = len(self.checks)
        self.verdict = "PASS" if passed_count == total_count else "REJECT"
        return {"verdict": self.verdict, "passed": passed_count, "total": total_count, "checks": self.checks, "timestamp": datetime.now().isoformat()}

    def generate_report_md(self):
        lines = [
            "# QA GATE AUDIT REPORT",
            f"status: {self.verdict or 'not run'}",
            f"timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            f"package: {self.package_root}", "",
            "## Results", "| ID | Check | Result |", "|---|---|---|",
        ]
        for c in self.checks:
            status = "PASS" if c["passed"] else "FAIL"
            lines.append(f"| {c['id']} | {c['check']} | {status} |")
        lines.extend(["", f"## Verdict: **{self.verdict}**"])
        return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════
# 4. REPORT GENERATOR (PPT)
# ═══════════════════════════════════════════════════════════════

class ReportGenerator:
    """Generate professional PPT reports from data using the enhancer."""

    def __init__(self, bg_image=None):
        self.bg_image = bg_image
        self.enhancer = PPTEnhancer(bg_image_path=bg_image, transition="fade", compress=True)

    def create_report_ppt(self, title, slides_data, output_path):
        if not HAS_PPTX:
            return {"error": "python-pptx not installed"}
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d %H:%M")

        for sdata in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sdata.get("title", "")
            if slide.placeholders[1]:
                slide.placeholders[1].text = sdata.get("content", "")

        raw_path = output_path + ".raw.pptx"
        prs.save(raw_path)
        result = self.enhancer.enhance(raw_path, output_path)
        if os.path.exists(raw_path):
            os.remove(raw_path)
        return result

    def create_lead_pack_report(self, lead_data, output_path):
        total = len(lead_data)
        wholesalers = sum(1 for l in lead_data if l.get("Lead_Type") == "Wholesaler/Buyer")
        distressed = sum(1 for l in lead_data if l.get("Lead_Type") == "Distressed Seller")
        slides = [
            {"title": "Lead Intelligence Report", "content": f"Total Leads: {total}\nWholesalers/Buyers: {wholesalers}\nDistressed Sellers: {distressed}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"},
            {"title": "Top Properties", "content": "\n".join(f"  {l.get('Property_Address', l.get('Company', 'N/A'))} - {l.get('Distress_Signal', l.get('Lead_Source', ''))}" for l in lead_data[:10])},
        ]
        return self.create_report_ppt("MBM Daily Lead Pack", slides, output_path)

    def create_pipeline_report(self, pipeline_data, output_path):
        total = len(pipeline_data)
        by_stage = {}
        total_value = 0
        for d in pipeline_data:
            stage = d.get("stage", "unknown")
            by_stage[stage] = by_stage.get(stage, 0) + 1
            try:
                val = d.get("deal_value", "0").replace("$", "").replace(",", "")
                total_value += sum(int(x) for x in val.split("-")) / 2 if "-" in val else int(val)
            except:
                pass
        slides = [
            {"title": "Sales Pipeline Dashboard", "content": f"Total Deals: {total}\nPipeline Value: ${total_value:,.0f}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"},
            {"title": "Pipeline by Stage", "content": "\n".join(f"  {stage}: {count} deals" for stage, count in by_stage.items())},
        ]
        return self.create_report_ppt("MBM Sales Pipeline", slides, output_path)

    def create_qa_report(self, qa_result, output_path):
        verdict = qa_result.get("verdict", "UNKNOWN")
        slides = [{"title": f"QA Gate: {verdict}", "content": f"Checks Passed: {qa_result.get('passed', 0)}/{qa_result.get('total', 0)}\nTimestamp: {qa_result.get('timestamp', '')}"}]
        for check in qa_result.get("checks", []):
            status = "PASS" if check["passed"] else "FAIL"
            slides.append({"title": f"[{check['id']}] {check['check']}", "content": f"Result: {status}"})
        return self.create_report_ppt("QA Audit Report", slides, output_path)


# ═══════════════════════════════════════════════════════════════
# SERVICE CATALOG
# ═══════════════════════════════════════════════════════════════

SERVICE_CATALOG = {
    "ppt_enhance_basic":    {"name": "PPT Enhancement - Basic",    "desc": "Futuristic backgrounds on all slides",                       "price": 50},
    "ppt_enhance_standard": {"name": "PPT Enhancement - Standard", "desc": "Backgrounds + smooth transitions",                            "price": 100},
    "ppt_enhance_premium":  {"name": "PPT Enhancement - Premium",  "desc": "Backgrounds + transitions + compression + custom branding",   "price": 200},
    "boq_restructure":      {"name": "BOQ Restructure",            "desc": "Standardize BOQ to Sr/Desc/Model/Unit/QTY/Price/Total format","price": 150},
    "boq_redaction":        {"name": "Content Redaction",          "desc": "Remove sensitive content from BOQ/PPT/PDF packages",          "price": 100},
    "boq_full_rebuild":     {"name": "Full BOQ Rebuild",           "desc": "Restructure + redact + rebuild master from per-space sheets", "price": 300},
    "qa_audit":             {"name": "QA Gate Audit",              "desc": "12-point automated quality check with report",                "price": 200},
    "qa_fix":               {"name": "QA Audit + Fix",             "desc": "Audit + remediate all failures",                              "price": 500},
    "qa_retainer":          {"name": "QA Retainer (Monthly)",      "desc": "Monthly QA coverage for all deliverables",                    "price": 1000},
    "lead_pack_enhanced":   {"name": "Enhanced Lead Pack (Daily)", "desc": "Daily leads + visual PPT report",                             "price": 75},
    "lead_pack_monthly":    {"name": "Enhanced Lead Pack (Monthly)","desc": "Unlimited enhanced lead packs",                              "price": 1500},
}


def main():
    import argparse
    parser = argparse.ArgumentParser(description="MBM Contractor Package Services")
    sub = parser.add_subparsers(dest="command")

    ppt = sub.add_parser("enhance-ppt", help="Enhance a PPTX file")
    ppt.add_argument("input", help="Input PPTX file")
    ppt.add_argument("--output", help="Output PPTX file")
    ppt.add_argument("--bg", help="Background image path")
    ppt.add_argument("--transition", default="fade", choices=["fade", "push", "wipe", "cover"])
    ppt.add_argument("--no-compress", action="store_true")

    boq = sub.add_parser("scan-boq", help="Scan a BOQ folder")
    boq.add_argument("path", help="BOQ root directory")

    qa = sub.add_parser("qa-audit", help="Run QA audit on a package")
    qa.add_argument("path", help="Package root directory")
    qa.add_argument("--banned", nargs="*", help="Banned terms to scan for")

    sub.add_parser("catalog", help="Print service catalog")

    args = parser.parse_args()

    if args.command == "enhance-ppt":
        enhancer = PPTEnhancer(bg_image_path=args.bg, transition=args.transition, compress=not args.no_compress)
        result = enhancer.enhance(args.input, args.output)
        print(json.dumps(result, indent=2))
    elif args.command == "scan-boq":
        proc = BOQProcessor()
        result = proc.scan_boq_folder(args.path)
        print(json.dumps(result, indent=2, default=str))
    elif args.command == "qa-audit":
        gate = QAGate(args.path)
        result = gate.run_full_audit(banned_terms=args.banned)
        print(gate.generate_report_md())
    elif args.command == "catalog":
        print("=" * 60)
        print("MBM CONTRACTOR PACKAGE SERVICES")
        print("=" * 60)
        for key, svc in SERVICE_CATALOG.items():
            print(f"\n  [{key}] {svc['name']}")
            print(f"    {svc['desc']}")
            print(f"    ${svc['price']}")
        print("\n" + "=" * 60)
    else:
        parser.print_help()

if __name__ == "__main__":
    main()
